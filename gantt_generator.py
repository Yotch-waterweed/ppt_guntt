import datetime
import calendar
from typing import List, Dict, Any
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class GanttChartGenerator:
    # 四半期の月マッピング（日本の会計年度ベース）
    QUARTER_MONTHS = {1: [4, 5, 6], 2: [7, 8, 9], 3: [10, 11, 12], 4: [1, 2, 3]}
    
    def __init__(self, data: List[Dict[str, Any]], start_month: str=None, end_month: str=None,
                 force_external_label_months: int=3,
                 calendar_mode: str="auto", bme_threshold: int=13):
        """
        Args:
            data: ガントチャートのデータ
            start_month: 表示開始月 (YYYY-MM)
            end_month: 表示終了月 (YYYY-MM)
            force_external_label_months: この月数以上で矢羽テキストを強制外出し
            calendar_mode: "auto", "bme", "monthly", "quarterly"
            bme_threshold: autoモード時、この月数以上でB/M/E→月単位に切替
        """
        self.data = data
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        self.force_external_label_months = force_external_label_months
        self.months = self._calculate_date_range(data, start_month, end_month)
        
        # カレンダーモード判定
        if calendar_mode == "auto":
            if len(self.months) < bme_threshold:
                self._cal_mode = "bme"
            else:
                self._cal_mode = "monthly"
        else:
            self._cal_mode = calendar_mode
        
        # 四半期モードの場合、四半期リストを生成
        if self._cal_mode == "quarterly":
            self.quarters = self._build_quarters()
            self.header_rows = 2
        elif self._cal_mode == "bme":
            self.quarters = []
            self.header_rows = 3
        else:
            self.quarters = []
            self.header_rows = 2
        
        self.layout = self._calculate_layout()
    
    def _build_quarters(self):
        """self.months から四半期リストを生成する"""
        quarters = []
        seen = set()
        for m in self.months:
            month = m["month"]
            # 会計年度のFY年とQ番号を計算
            if month >= 4:
                fy_year = m["year"]
                q = (month - 4) // 3 + 1  # 4-6→Q1, 7-9→Q2, 10-12→Q3
            else:
                fy_year = m["year"] - 1
                q = 4  # 1-3月→Q4
            
            key = (fy_year, q)
            if key not in seen:
                seen.add(key)
                # この四半期に含まれる月の総日数を計算
                q_months = self.QUARTER_MONTHS[q]
                cal_year_for_q = fy_year if q <= 3 else fy_year + 1
                total_days = 0
                for qm in q_months:
                    y = cal_year_for_q if qm >= 4 else cal_year_for_q
                    if q == 4:
                        y = fy_year + 1  # Q4の1-3月は翌暦年
                    elif qm >= 4:
                        y = fy_year
                    else:
                        y = fy_year + 1
                    total_days += calendar.monthrange(y, qm)[1]
                
                quarters.append({
                    "fy_year": fy_year,
                    "quarter": q,
                    "months": q_months,
                    "total_days": total_days,
                    "label": f"Q{q}",
                })
        return quarters
    
    def _calculate_layout(self):
        """スライドの使用可能高さに基づき、各要素のサイズを自動計算する"""
        base = {
            "header_h": Pt(25),
            "bar_h": Pt(15), "ms_size": Pt(15), "ms_clearance": Pt(6),
            "chevron_h": Pt(20), "chevron_gap": Pt(4),
            "section_gap": Pt(8), "row_margin": Pt(6),
        }
        minimum = {
            "bar_h": Pt(8), "ms_size": Pt(10), "ms_clearance": Pt(3),
            "chevron_h": Pt(12), "chevron_gap": Pt(2),
            "section_gap": Pt(4), "row_margin": Pt(3),
        }
        
        available_h = Inches(7.5) - Inches(1.0) - base["header_h"] * self.header_rows
        
        total_base_h = 0
        for item in self.data:
            has_items = len(item.get("items", [])) > 0
            num_tasks = len(item.get("tasks", []))
            upper_h = (base["bar_h"] + base["ms_size"] + base["ms_clearance"]) if has_items else 0
            lower_h = (num_tasks * (base["chevron_h"] + base["chevron_gap"]) - base["chevron_gap"]) if num_tasks > 0 else 0
            row_h = base["row_margin"] * 2
            if upper_h > 0 and lower_h > 0:
                row_h += upper_h + base["section_gap"] + lower_h
            elif upper_h > 0:
                row_h += upper_h
            elif lower_h > 0:
                row_h += lower_h
            total_base_h += max(int(row_h), Pt(50))
        
        scale = min(1.0, available_h / total_base_h) if total_base_h > 0 else 1.0
        
        layout = {"header_h": base["header_h"], "scale": scale}
        for key in minimum:
            layout[key] = max(int(base[key] * scale), minimum[key])
        
        # カレンダー列数
        if self._cal_mode == "bme":
            num_cal_cols = len(self.months) * 3
        elif self._cal_mode == "quarterly":
            num_cal_cols = len(self.quarters)
        else:
            num_cal_cols = len(self.months)
        
        num_months = len(self.months)
        if num_months <= 6:
            layout["header_font"] = Pt(12)
            layout["bme_font"] = Pt(10)
            layout["task_name_width"] = Inches(2.5)
            layout["task_name_font"] = Pt(12)
        elif num_months <= 12:
            layout["header_font"] = Pt(9)
            layout["bme_font"] = Pt(7)
            layout["task_name_width"] = Inches(2.0)
            layout["task_name_font"] = Pt(10)
        else:
            layout["header_font"] = Pt(8)
            layout["bme_font"] = Pt(6)
            layout["task_name_width"] = Inches(1.8)
            layout["task_name_font"] = Pt(9)
        
        layout["num_cal_cols"] = num_cal_cols
        return layout
        
    def _calculate_date_range(self, data, start_month, end_month):
        """全体の開始・終了月を計算し、カレンダー用データの配列を生成する"""
        min_date = datetime.date.max
        max_date = datetime.date.min
        
        for item in data:
            if item.get("type") == "project":
                for p_item in item.get("items", []):
                    if p_item.get("type") == "milestone" and "date" in p_item:
                        d = datetime.date.fromisoformat(p_item["date"])
                        min_date, max_date = min(min_date, d), max(max_date, d)
                    elif p_item.get("type") == "period" and "start" in p_item and "end" in p_item:
                        d_s = datetime.date.fromisoformat(p_item["start"])
                        d_e = datetime.date.fromisoformat(p_item["end"])
                        min_date, max_date = min(min_date, d_s), max(max_date, d_e)
                for t_item in item.get("tasks", []):
                    if "start" in t_item and "end" in t_item:
                        d_s = datetime.date.fromisoformat(t_item["start"])
                        d_e = datetime.date.fromisoformat(t_item["end"])
                        min_date, max_date = min(min_date, d_s), max(max_date, d_e)
        
        if start_month: min_date = datetime.date.fromisoformat(start_month + "-01")
        if end_month:
            y, m = map(int, end_month.split('-'))
            max_date = datetime.date(y, m, calendar.monthrange(y, m)[1])
            
        start_cal = datetime.date(min_date.year, min_date.month, 1)
        end_cal = datetime.date(max_date.year, max_date.month,
                                calendar.monthrange(max_date.year, max_date.month)[1])
        
        months = []
        current = start_cal
        while current <= end_cal:
            months.append({
                "year": current.year,
                "month": current.month,
                "days": calendar.monthrange(current.year, current.month)[1]
            })
            current += relativedelta(months=1)
        return months
        
    def _set_cell_style(self, cell, text, bg_color=None, is_bold=False, font_size=None):
        """セルのテキストと背景色を設定する共通処理"""
        if font_size is None:
            font_size = self.layout["header_font"]
        cell.text = text
        cell.margin_left = Pt(1)
        cell.margin_right = Pt(1)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.bold = is_bold
                if bg_color:
                    run.font.color.rgb = RGBColor(255, 255, 255)
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    def _create_calendar_headers(self, table):
        """カレンダーのヘッダーを描画する"""
        header_bg = RGBColor(89, 89, 89)
        
        # 左上の空白部分を結合
        table.cell(0, 0).merge(table.cell(self.header_rows - 1, 0))
        self._set_cell_style(table.cell(0, 0), "タスク名 / 担当", header_bg, True)
        
        if self._cal_mode == "quarterly":
            self._create_quarterly_headers(table, header_bg)
        elif self._cal_mode == "bme":
            self._create_bme_headers(table, header_bg)
        else:
            self._create_monthly_headers(table, header_bg)
    
    def _create_quarterly_headers(self, table, header_bg):
        """四半期モードのヘッダー: FY行 + Q行"""
        # FY行の結合範囲を計算
        fy_ranges = []
        current_fy = None
        start_idx = 0
        for i, q in enumerate(self.quarters):
            if current_fy != q["fy_year"]:
                if current_fy is not None:
                    fy_ranges.append((current_fy, start_idx, i - 1))
                current_fy = q["fy_year"]
                start_idx = i
        fy_ranges.append((current_fy, start_idx, len(self.quarters) - 1))
        
        # 1段目: FY
        for fy_year, s_idx, e_idx in fy_ranges:
            start_col = 1 + s_idx
            end_col = 1 + e_idx
            cell = table.cell(0, start_col)
            if start_col != end_col:
                cell.merge(table.cell(0, end_col))
            self._set_cell_style(cell, f"FY{str(fy_year)[-2:]}", header_bg, True)
        
        # 2段目: Q1, Q2, Q3, Q4
        for i, q in enumerate(self.quarters):
            col = 1 + i
            self._set_cell_style(table.cell(1, col), q["label"], header_bg, True)
    
    def _create_bme_headers(self, table, header_bg):
        """B/M/Eモードのヘッダー"""
        year_ranges = self._calc_year_ranges(3)
        for year, s_idx, e_idx in year_ranges:
            start_col = 1 + s_idx * 3
            end_col = 1 + e_idx * 3 + 2
            cell = table.cell(0, start_col)
            if start_col != end_col:
                cell.merge(table.cell(0, end_col))
            self._set_cell_style(cell, f"FY{str(year)[-2:]}", header_bg, True)
        
        for i, m in enumerate(self.months):
            start_col = 1 + i * 3
            cell_month = table.cell(1, start_col)
            cell_month.merge(table.cell(1, start_col + 2))
            self._set_cell_style(cell_month, str(m['month']), header_bg, True)
            self._set_cell_style(table.cell(2, start_col), "B", header_bg, font_size=self.layout["bme_font"])
            self._set_cell_style(table.cell(2, start_col + 1), "M", header_bg, font_size=self.layout["bme_font"])
            self._set_cell_style(table.cell(2, start_col + 2), "E", header_bg, font_size=self.layout["bme_font"])
    
    def _create_monthly_headers(self, table, header_bg):
        """月単位モードのヘッダー"""
        year_ranges = self._calc_year_ranges(1)
        for year, s_idx, e_idx in year_ranges:
            start_col = 1 + s_idx
            end_col = 1 + e_idx
            cell = table.cell(0, start_col)
            if start_col != end_col:
                cell.merge(table.cell(0, end_col))
            self._set_cell_style(cell, f"FY{str(year)[-2:]}", header_bg, True)
        
        for i, m in enumerate(self.months):
            col = 1 + i
            self._set_cell_style(table.cell(1, col), str(m['month']), header_bg, True)
    
    def _calc_year_ranges(self, cols_per_month):
        """年ごとの列結合範囲を計算する"""
        year_ranges = []
        current_year = None
        start_idx = 0
        for i, m in enumerate(self.months):
            if current_year != m["year"]:
                if current_year is not None:
                    year_ranges.append((current_year, start_idx, i - 1))
                current_year = m["year"]
                start_idx = i
        year_ranges.append((current_year, start_idx, len(self.months) - 1))
        return year_ranges

    def _populate_task_names(self, table):
        """左の列にプロジェクト名を入力する"""
        row_idx = self.header_rows
        for item in self.data:
            cell = table.cell(row_idx, 0)
            cell.text = item.get("name", "")
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.bold = True
                run.font.size = self.layout["task_name_font"]
            row_idx += 1

    def _draw_table_and_background(self, slide):
        """全体のテーブルとグリッドを描画する"""
        num_cal_cols = self.layout["num_cal_cols"]
        num_cols = 1 + num_cal_cols
        num_rows = self.header_rows + len(self.data)
        
        left, top = Inches(0.5), Inches(0.5)
        width, height = Inches(12.333), Inches(6.5)
        
        shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = shape.table
        
        for i in range(self.header_rows):
            table.rows[i].height = Pt(25)
        
        L = self.layout
        for i, item in enumerate(self.data):
            has_items = len(item.get("items", [])) > 0
            num_tasks = len(item.get("tasks", []))
            upper_h = (L["bar_h"] + L["ms_size"] + L["ms_clearance"]) if has_items else 0
            lower_h = (num_tasks * (L["chevron_h"] + L["chevron_gap"]) - L["chevron_gap"]) if num_tasks > 0 else 0
            total_h = L["row_margin"] * 2
            if upper_h > 0 and lower_h > 0:
                total_h += upper_h + L["section_gap"] + lower_h
            elif upper_h > 0:
                total_h += upper_h
            elif lower_h > 0:
                total_h += lower_h
            min_row_h = int(Pt(50) * L["scale"]) if L["scale"] < 1.0 else Pt(50)
            table.rows[self.header_rows + i].height = max(int(total_h), min_row_h)
        
        task_col_width = L["task_name_width"]
        table.columns[0].width = int(task_col_width)
        cal_col_width = (width - task_col_width) / num_cal_cols
        for i in range(1, num_cols):
            table.columns[i].width = int(cal_col_width)
            
        self._create_calendar_headers(table)
        self._populate_task_names(table)
        return shape
        
    def _get_x_coordinate(self, table_shape, date_str: str) -> int:
        """日付文字列から対応するテーブル上のX座標を計算する"""
        table = table_shape.table
        d = datetime.date.fromisoformat(date_str)
        
        if self._cal_mode == "quarterly":
            return self._get_x_quarterly(table_shape, d)
        
        # --- B/M/E or Monthly モード ---
        month_idx = -1
        for i, m in enumerate(self.months):
            if m["year"] == d.year and m["month"] == d.month:
                month_idx = i
                break
        
        if month_idx == -1:
            if d < datetime.date(self.months[0]["year"], self.months[0]["month"], 1):
                return table_shape.left + table.columns[0].width
            return table_shape.left + table_shape.width
        
        day = d.day
        days_in_month = self.months[month_idx]["days"]
        
        if self._cal_mode == "bme":
            if day <= 10:
                col_offset, fraction = 0, (day - 1) / 10.0
            elif day <= 20:
                col_offset, fraction = 1, (day - 11) / 10.0
            else:
                col_offset, fraction = 2, (day - 21) / max(1.0, float(days_in_month - 20))
            col_idx = 1 + month_idx * 3 + col_offset
        else:
            fraction = (day - 1) / max(1.0, float(days_in_month - 1))
            col_idx = 1 + month_idx
        
        x = table_shape.left
        for i in range(col_idx):
            x += table.columns[i].width
        x += int(table.columns[col_idx].width * fraction)
        return x
    
    def _get_x_quarterly(self, table_shape, d):
        """四半期モード用のX座標計算"""
        table = table_shape.table
        month = d.month
        
        # この日付が属する四半期を特定
        if month >= 4:
            fy_year = d.year
            q = (month - 4) // 3 + 1
        else:
            fy_year = d.year - 1
            q = 4
        
        # 四半期リスト内のインデックスを探す
        q_idx = -1
        for i, qinfo in enumerate(self.quarters):
            if qinfo["fy_year"] == fy_year and qinfo["quarter"] == q:
                q_idx = i
                break
        
        if q_idx == -1:
            # 範囲外
            first_q = self.quarters[0]
            first_month = first_q["months"][0]
            first_year = first_q["fy_year"] if first_month >= 4 else first_q["fy_year"] + 1
            if d < datetime.date(first_year, first_month, 1):
                return table_shape.left + table.columns[0].width
            return table_shape.left + table_shape.width
        
        # 四半期内での経過日数を計算
        qinfo = self.quarters[q_idx]
        q_months = qinfo["months"]
        # 四半期の開始日を計算
        first_m = q_months[0]
        if first_m >= 4:
            q_start_year = fy_year
        else:
            q_start_year = fy_year + 1
        q_start = datetime.date(q_start_year, first_m, 1)
        
        elapsed_days = (d - q_start).days
        fraction = max(0.0, min(1.0, elapsed_days / max(1, qinfo["total_days"])))
        
        col_idx = 1 + q_idx
        x = table_shape.left
        for i in range(col_idx):
            x += table.columns[i].width
        x += int(table.columns[col_idx].width * fraction)
        return x

    def _draw_shapes(self, slide, table_shape):
        """ガントバーやマイルストーンの図形を描画する"""
        from pptx.enum.shapes import MSO_SHAPE
        
        table = table_shape.table
        color_project = RGBColor(100, 100, 200)
        color_task = RGBColor(130, 130, 130)
        color_milestone = RGBColor(220, 80, 80)
        color_text_dark = RGBColor(50, 50, 50)
        slide_right = self.prs.slide_width - Inches(0.1)
        
        L = self.layout
        bar_h, ms_size, ms_clearance = L["bar_h"], L["ms_size"], L["ms_clearance"]
        chevron_h, chevron_gap = L["chevron_h"], L["chevron_gap"]
        section_gap, row_margin = L["section_gap"], L["row_margin"]
        
        for i, item in enumerate(self.data):
            row_idx = self.header_rows + i
            row_y = table_shape.top + sum(table.rows[r].height for r in range(row_idx))
            row_h = table.rows[row_idx].height
            
            has_items = len(item.get("items", [])) > 0
            tasks = item.get("tasks", [])
            has_tasks = len(tasks) > 0
            
            if has_items and has_tasks:
                bar_y = int(row_y + row_margin + ms_size + ms_clearance)
                tasks_start_y = int(bar_y + bar_h + section_gap)
            elif has_items:
                bar_y = int(row_y + row_margin + ms_size + ms_clearance)
                tasks_start_y = 0
            else:
                bar_y = 0
                total_tasks_h = len(tasks) * chevron_h + max(0, len(tasks) - 1) * chevron_gap
                tasks_start_y = int(row_y + (row_h - total_tasks_h) / 2)
            
            if has_items:
                for p_item in item.get("items", []):
                    if p_item.get("type") == "period" and "start" in p_item and "end" in p_item:
                        x_start = self._get_x_coordinate(table_shape, p_item["start"])
                        x_end = self._get_x_coordinate(table_shape, p_item["end"])
                        w = max(x_end - x_start, Inches(0.05))
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, bar_y, w, bar_h)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = color_project
                        shape.line.color.rgb = color_project
                        
                        text = p_item.get("name", "")
                        if text:
                            if len(text) * Pt(12) < w * 0.9:
                                shape.text = text
                                shape.text_frame.word_wrap = False
                                shape.text_frame.paragraphs[0].font.size = Pt(9)
                                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            else:
                                tx_x = int(x_start + w + Pt(1))
                                tx_w = min(Inches(2), max(int(slide_right - tx_x), Pt(10)))
                                if tx_x < slide_right:
                                    self._add_label(slide, tx_x, bar_y, tx_w, bar_h, text, color_text_dark)
                
                for p_item in item.get("items", []):
                    if p_item.get("type") == "milestone" and "date" in p_item:
                        x = self._get_x_coordinate(table_shape, p_item["date"])
                        ms_y = int(bar_y - ms_size - ms_clearance)
                        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, int(x - ms_size/2), ms_y, ms_size, ms_size)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = color_milestone
                        shape.line.color.rgb = color_milestone
                        shape.line.width = Pt(1)
                        
                        text = p_item.get("name", "")
                        if text:
                            tx_x = int(x + ms_size/2 + Pt(1))
                            tx_w = min(Inches(1.5), max(int(slide_right - tx_x), Pt(10)))
                            if tx_x < slide_right:
                                self._add_label(slide, tx_x, ms_y, tx_w, ms_size, text, color_text_dark)
            
            if has_tasks:
                for t_idx, t_item in enumerate(tasks):
                    if "start" in t_item and "end" in t_item:
                        x_start = self._get_x_coordinate(table_shape, t_item["start"])
                        x_end = self._get_x_coordinate(table_shape, t_item["end"])
                        w = max(x_end - x_start, Inches(0.2))
                        chevron_y = int(tasks_start_y + t_idx * (chevron_h + chevron_gap))
                        
                        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x_start, chevron_y, w, chevron_h)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = color_task
                        shape.line.color.rgb = color_task
                        
                        text = t_item.get("name", "")
                        if text:
                            force_ext = len(self.months) >= self.force_external_label_months
                            if not force_ext and len(text) * Pt(7) < w * 0.85:
                                shape.text = text
                                shape.text_frame.word_wrap = False
                                shape.text_frame.paragraphs[0].font.size = Pt(9)
                                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            else:
                                tx_x = int(x_start + w + Pt(0))
                                tx_w = min(Inches(2), max(int(slide_right - tx_x), Pt(10)))
                                if tx_x < slide_right:
                                    self._add_label(slide, tx_x, chevron_y, tx_w, chevron_h, text, color_text_dark)

    def _add_label(self, slide, x, y, w, h, text, color):
        """テキストボックスを追加するヘルパー"""
        txBox = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = color

    def generate(self, output_path: str):
        """PowerPointファイルを生成する"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        table_shape = self._draw_table_and_background(slide)
        self._draw_shapes(slide, table_shape)
        self.prs.save(output_path)
        print(f"[{output_path}] が生成されました。")

if __name__ == "__main__":
    mock_data = [
        {
            "name": "TOBE設計",
            "type": "project",
            "items": [
                {"name": "キックオフ", "type": "milestone", "date": "2025-04-01"},
                {"name": "前半フェーズ", "type": "period", "start": "2025-04-01", "end": "2025-05-15"},
                {"name": "中間レビュー", "type": "milestone", "date": "2025-05-15"},
                {"name": "後半フェーズ", "type": "period", "start": "2025-05-20", "end": "2025-08-15"}
            ],
            "tasks": []
        },
        {
            "name": "業務要件定義",
            "type": "project",
            "items": [
                {"name": "業務要件定義", "type": "period", "start": "2025-04-01", "end": "2025-04-30"}
            ],
            "tasks": [
                {"name": "現状業務フロー整理", "start": "2025-04-01", "end": "2025-04-10"},
                {"name": "課題・要望ヒアリング", "start": "2025-04-08", "end": "2025-04-18"},
                {"name": "業務要件書作成・レビュー", "start": "2025-04-16", "end": "2025-04-30"}
            ]
        },
        {
            "name": "システム要件定義",
            "type": "project",
            "items": [
                {"name": "システム要件定義", "type": "period", "start": "2025-04-21", "end": "2025-06-15"},
                {"name": "要件FIX", "type": "milestone", "date": "2025-06-15"}
            ],
            "tasks": [
                {"name": "機能要件洗い出し", "start": "2025-04-21", "end": "2025-05-05"},
                {"name": "非機能要件定義", "start": "2025-04-28", "end": "2025-05-10"},
                {"name": "外部IF仕様検討", "start": "2025-05-06", "end": "2025-05-20"},
                {"name": "データモデル設計", "start": "2025-05-12", "end": "2025-05-31"},
                {"name": "要件定義書作成・承認", "start": "2025-05-26", "end": "2025-06-15"}
            ]
        }
    ]
    
    generator = GanttChartGenerator(mock_data)
    generator.generate("test_gantt.pptx")
